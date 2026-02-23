#!/usr/bin/env python3
"""
Extract storyboard panel images from PPTX files.

Reads image placement coordinates and text boxes directly from the
PowerPoint structure — no rasterization or computer vision needed.

Usage:
  python3 pptx_panels.py <pptx_path> <slide_number>

Returns JSON with same format as pdf_panels.py:
  {
    "count": 6,
    "mode": "pptx_structure",
    "panels": [
      { "x": 100, "y": 200, "width": 300, "height": 200, "image": "<base64>", "caption": "..." },
      ...
    ]
  }
"""
import json
import sys
import base64
import io

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN


def emu_to_pt(emu):
    """Convert EMU (English Metric Units) to points."""
    return emu / 12700


def extract_panels(pptx_path, slide_num, min_w_pt=80, min_h_pt=50):
    """
    Extract panel images from a PPTX slide.
    
    1. Find all PICTURE shapes on the slide
    2. Filter out small logos/icons by minimum size
    3. Sort in reading order (top-to-bottom, left-to-right)
    4. Extract image bytes
    5. Find associated text boxes (captions below/near each image)
    6. Return with formatting markup
    """
    prs = Presentation(pptx_path)
    
    if slide_num < 1 or slide_num > len(prs.slides):
        return {"count": 0, "mode": "pptx_structure", "panels": [],
                "error": f"Slide {slide_num} out of range (1-{len(prs.slides)})"}
    
    slide = prs.slides[slide_num - 1]
    slide_w_pt = emu_to_pt(prs.slide_width)
    slide_h_pt = emu_to_pt(prs.slide_height)
    
    # Collect all picture shapes with position info
    pictures = []
    text_shapes = []
    title_text = None
    
    for shape in slide.shapes:
        left_pt = emu_to_pt(shape.left) if shape.left is not None else 0
        top_pt = emu_to_pt(shape.top) if shape.top is not None else 0
        w_pt = emu_to_pt(shape.width) if shape.width is not None else 0
        h_pt = emu_to_pt(shape.height) if shape.height is not None else 0
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
            shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_text_frame is False and hasattr(shape, 'image')
        ):
            try:
                img_blob = shape.image.blob
                content_type = shape.image.content_type
                if w_pt >= min_w_pt and h_pt >= min_h_pt:
                    pictures.append({
                        "left": left_pt, "top": top_pt,
                        "width": w_pt, "height": h_pt,
                        "blob": img_blob,
                        "content_type": content_type
                    })
                else:
                    print(f"[PPTX] Slide {slide_num}: skipped small image {w_pt:.0f}x{h_pt:.0f}pt",
                          file=sys.stderr, flush=True)
            except Exception:
                pass  # Shape looks like picture but has no image data
        
        if shape.has_text_frame:
            text_shapes.append({
                "left": left_pt, "top": top_pt,
                "width": w_pt, "height": h_pt,
                "text_frame": shape.text_frame
            })
            
            # Check if this is a title (near top of slide, wide, short)
            if top_pt < slide_h_pt * 0.20 and w_pt > slide_w_pt * 0.3:
                raw = shape.text_frame.text.strip()
                if raw and not title_text:
                    title_text = raw
    
    if not pictures:
        print(f"[PPTX] Slide {slide_num}: no panel images found", file=sys.stderr, flush=True)
        return {"count": 0, "mode": "pptx_structure", "panels": []}
    
    # Sort reading order: group by rows, then left-to-right within each row
    pictures.sort(key=lambda p: p["top"])
    row_threshold = 20  # points
    rows = []
    current_row = [pictures[0]]
    for p in pictures[1:]:
        if abs(p["top"] - current_row[0]["top"]) < row_threshold:
            current_row.append(p)
        else:
            rows.append(sorted(current_row, key=lambda p: p["left"]))
            current_row = [p]
    rows.append(sorted(current_row, key=lambda p: p["left"]))
    pictures = [p for row in rows for p in row]
    
    # Build result panels
    result_panels = []
    for pic in pictures:
        # Convert image to base64 JPEG
        blob = pic["blob"]
        ct = pic["content_type"]
        
        # If already JPEG, use as-is. Otherwise convert via PIL if available.
        if ct in ("image/jpeg", "image/jpg"):
            b64 = base64.b64encode(blob).decode("ascii")
        else:
            try:
                from PIL import Image
                img = Image.open(io.BytesIO(blob))
                if img.mode in ("RGBA", "P", "LA"):
                    img = img.convert("RGB")
                buf = io.BytesIO()
                img.save(buf, format="JPEG", quality=85)
                b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            except ImportError:
                # No PIL — just base64 the raw bytes, let Node handle conversion
                b64 = base64.b64encode(blob).decode("ascii")
        
        # Find caption: text box whose top is near this image's bottom,
        # with horizontal overlap
        img_bottom = pic["top"] + pic["height"]
        img_left = pic["left"]
        img_right = pic["left"] + pic["width"]
        
        caption_text = ""
        best_dist = 999
        for ts in text_shapes:
            ts_top = ts["top"]
            y_dist = ts_top - img_bottom
            ts_right = ts["left"] + ts["width"]
            h_overlap = min(ts_right, img_right) - max(ts["left"], img_left)
            
            if -10 <= y_dist < 60 and h_overlap > 30 and y_dist < best_dist:
                # Extract formatted text
                lines = []
                for para in ts["text_frame"].paragraphs:
                    line_parts = []
                    for run in para.runs:
                        text = run.text
                        if not text.strip():
                            line_parts.append(text)
                            continue
                        bold = run.font.bold
                        italic = run.font.italic
                        if bold and italic:
                            line_parts.append(f"<b><i>{text}</i></b>")
                        elif bold:
                            line_parts.append(f"<b>{text}</b>")
                        elif italic:
                            line_parts.append(f"<i>{text}</i>")
                        else:
                            line_parts.append(text)
                    line = "".join(line_parts).strip()
                    if line:
                        lines.append(line)
                
                candidate = "\n".join(lines)
                if candidate:
                    caption_text = candidate
                    best_dist = y_dist
        
        result_panels.append({
            "x": round(pic["left"], 1),
            "y": round(pic["top"], 1),
            "width": round(pic["width"], 1),
            "height": round(pic["height"], 1),
            "image": b64,
            "caption": caption_text
        })
    
    print(f"[PPTX] Slide {slide_num}: {len(result_panels)} panels extracted from PPTX structure",
          file=sys.stderr, flush=True)
    
    return {
        "count": len(result_panels),
        "mode": "pptx_structure",
        "panels": result_panels,
        "title": title_text
    }


def get_slide_count(pptx_path):
    """Return the number of slides in a PPTX file."""
    prs = Presentation(pptx_path)
    return len(prs.slides)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        # If only path provided, return slide count
        if len(sys.argv) == 2:
            try:
                count = get_slide_count(sys.argv[1])
                print(json.dumps({"slideCount": count}))
            except Exception as e:
                print(json.dumps({"error": str(e)}))
                sys.exit(1)
        else:
            print(json.dumps({"error": "Usage: pptx_panels.py <pptx_path> [slide_number]"}))
            sys.exit(1)
    else:
        pptx_path = sys.argv[1]
        slide_num = int(sys.argv[2])
        
        try:
            result = extract_panels(pptx_path, slide_num)
            print(json.dumps(result))
        except Exception as e:
            print(json.dumps({"error": str(e)}))
            sys.exit(1)
